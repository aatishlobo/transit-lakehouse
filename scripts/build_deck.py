"""Generate the final presentation deck.

Built to the published rubric, which awards 5 points each for:

  1. Clear problem statement -- problem, target user, useful result, why it matters
  2. Data and event contract -- source, Kafka event, key fields, ONE representative record
  3. Architecture chart and design choice -- end-to-end path, ONE concrete example, WHY
  4. Lessons learned and future work -- one evidence-based lesson, one realistic
     next step, AND staying within time

Slot: Thursday 13 August, 6:16-6:24pm. Two-person team = 7 minutes presenting
plus 1 minute Q&A. Both speakers must talk.

Timing is the part most decks get wrong, so each slide carries a budget in its
speaker notes and the totals are asserted at the end of this script. 12 slides
in 7 minutes is ~35s average -- the deck is deliberately built so most slides
are a single idea that can be delivered in one breath.

Regenerate:  python scripts/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palette -- restrained, high contrast, projector-safe.
# Light background because a classroom projector washes out dark themes.
# ---------------------------------------------------------------------------
INK = RGBColor(0x1A, 0x1F, 0x2B)      # near-black body text
MUTED = RGBColor(0x5B, 0x66, 0x7A)    # secondary text
ACCENT = RGBColor(0x0B, 0x6E, 0x99)   # teal-blue, the one accent
WARN = RGBColor(0xB4, 0x53, 0x09)     # amber, used ONLY for the key finding
BG = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF5, 0xF8)
LINE = RGBColor(0xD5, 0xDD, 0xE5)

W, H = Inches(13.333), Inches(7.5)    # 16:9
MARGIN = Inches(0.7)

TOTAL_BUDGET_S = 7 * 60


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, text, size=18, bold=False, color=INK,
        align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        # inline **bold** segments
        parts = line.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            r.font.size = Pt(size)
            r.font.bold = bold or (j % 2 == 1)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def title(slide, text, sub=None, speaker=None):
    txt(slide, MARGIN, Inches(0.45), W - 2 * MARGIN, Inches(0.8),
        text, size=32, bold=True, color=INK)
    y = Inches(1.15)
    if sub:
        txt(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.5), sub,
            size=17, color=MUTED)
        y = Inches(1.72)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y,
                                  Inches(1.6), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False
    if speaker:
        tb = txt(slide, W - MARGIN - Inches(2.4), Inches(0.5), Inches(2.4),
                 Inches(0.3), speaker, size=12, color=MUTED,
                 align=PP_ALIGN.RIGHT)
    return Inches(2.25)


def panel(slide, x, y, w, h, fill=PANEL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    box.adjustments[0] = 0.04
    return box


def mono(slide, x, y, w, h, code, size=13, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        p.line_spacing = 1.05
        parts = line.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            r.font.size = Pt(size)
            r.font.name = "Consolas"
            r.font.bold = j % 2 == 1
            r.font.color.rgb = WARN if j % 2 == 1 else color
    return tb


def notes(slide, seconds, text):
    slide.notes_slide.notes_text_frame.text = f"[{seconds}s]  {text}"
    return seconds


def flow_box(slide, x, y, w, h, label, sub=None, fill=PANEL, edge=LINE,
             bold_label=True):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = edge
    b.line.width = Pt(1.25)
    b.shadow.inherit = False
    b.adjustments[0] = 0.12
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(13)
    r.font.bold = bold_label
    r.font.color.rgb = INK
    r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED
        r2.font.name = "Calibri"
    return b


def arrow(slide, x, y, w, color=ACCENT):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.16))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# ---------------------------------------------------------------------------
def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    budget = []

    # ---------------- 1. Title -------------------------------------------
    s = add_slide(prs)
    txt(s, MARGIN, Inches(2.35), W - 2 * MARGIN, Inches(1.0),
        "Real-Time Transit Reliability Lakehouse", size=42, bold=True)
    txt(s, MARGIN, Inches(3.35), W - 2 * MARGIN, Inches(0.7),
        "Deriving arrival times that no agency publishes —\n"
        "and grading the predictions riders actually see",
        size=20, color=MUTED, line_spacing=1.3)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.45),
                           Inches(2.0), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT
    r.line.fill.background(); r.shadow.inherit = False
    txt(s, MARGIN, Inches(4.85), W - 2 * MARGIN, Inches(0.9),
        "Borna Karimi  ·  Aatish Lobo\nMSDS 682 — Data Stream Processing",
        size=17, color=INK, line_spacing=1.35)
    txt(s, MARGIN, Inches(6.6), W - 2 * MARGIN, Inches(0.3),
        "Data provided by 511.org (Metropolitan Transportation Commission)",
        size=11, color=MUTED)
    budget.append(notes(s, 10,
        "BORNA. 'We measure how late Bay Area transit actually is — by deriving "
        "arrival times that no agency publishes.' Move fast, the next slide is "
        "the hook."))

    # ---------------- 2. The problem -------------------------------------
    s = add_slide(prs)
    y = title(s, "Agencies never publish when a bus arrives",
              speaker="Borna")
    cw = (W - 2 * MARGIN - Inches(0.35)) / 2
    panel(s, MARGIN, y, cw, Inches(1.55))
    txt(s, MARGIN + Inches(0.25), y + Inches(0.22), cw - Inches(0.5), Inches(1.1),
        "TripUpdates — **predictions**\n“will reach stop 22 at 15:07”",
        size=16, line_spacing=1.35)
    panel(s, MARGIN + cw + Inches(0.35), y, cw, Inches(1.55))
    txt(s, MARGIN + cw + Inches(0.6), y + Inches(0.22), cw - Inches(0.5),
        Inches(1.1),
        "VehiclePositions — **GPS pings**\n“bus 4821 is here, now”",
        size=16, line_spacing=1.35)

    txt(s, MARGIN, y + Inches(1.95), W - 2 * MARGIN, Inches(0.6),
        "Neither ever says:  “bus 4821 **arrived** at stop 22 at 15:09.”",
        size=22, color=INK)
    txt(s, MARGIN, y + Inches(2.6), W - 2 * MARGIN, Inches(1.2),
        "That fact has to be **derived**.\n"
        "delay = actual − scheduled,  and scheduled is exact —\n"
        "so **every error in our derivation lands straight in the answer.**",
        size=18, color=INK, line_spacing=1.45)
    budget.append(notes(s, 50,
        "BORNA. The core surprise. Two feeds, neither states an arrival. "
        "Emphasise the last line: scheduled is exact, so derivation error goes "
        "directly into the delay number and nothing downstream can fix it."))

    # ---------------- 4. Data + contract ---------------------------------
    s = add_slide(prs)
    y = title(s, "Data source and the Kafka event contract",
              sub="511.org regional feed · GTFS-Realtime protobuf · 28 operators · 60 req/hour",
              speaker="Borna")
    panel(s, MARGIN, y - Inches(0.1), Inches(7.2), Inches(3.5))
    mono(s, MARGIN + Inches(0.3), y + Inches(0.12), Inches(6.7), Inches(3.2),
         """{
  "envelope": {
    "feed_header_ts":  "2026-08-06T21:52:47+00:00",
    "ingest_ts":       "2026-08-08T00:10:02+00:00",
    "poll_interval_s": 120,
    "contract_version":"1.0.0"
  },
  "trip": {
    "trip_id":      "SF:12053041_M11",
    "route_id":     "SF:1",
    "service_date": "20260806"
  },
  "stop_sequence":       47,
  "arrival_time_epoch":  1786053176,
  "arrival_delay_s":     140,
  "arrival_uncertainty": **null**
}""", size=12.5)

    xr = MARGIN + Inches(7.6)
    wr = W - MARGIN - xr
    txt(s, xr, y, wr, Inches(0.4), "Key fields", size=16, bold=True, color=ACCENT)
    txt(s, xr, y + Inches(0.45), wr, Inches(2.6),
        "**Kafka key:** `service_date:trip_id`\nkeeps one trip in one partition,\nin order\n\n"
        "**Grain:** service_date + trip_id +\nstop_sequence — never stop_id\n\n"
        "**poll_interval_s** travels per row:\nthe error bar on every arrival",
        size=14, line_spacing=1.35)
    txt(s, xr, y + Inches(2.95), wr, Inches(0.8),
        "**null ≠ 0** — and that\ndistinction is the whole project",
        size=15, color=WARN, line_spacing=1.3)
    budget.append(notes(s, 58,
        "BORNA. Rubric slide — one representative record. Point at "
        "arrival_uncertainty: null. Say the key is service_date:trip_id because "
        "trip IDs repeat daily. Then hand to the next slide as the punchline."))

    # ---------------- 5. The 44% finding ---------------------------------
    s = add_slide(prs)
    y = title(s, "Absent is not zero", speaker="Borna")
    txt(s, MARGIN, y - Inches(0.15), W - 2 * MARGIN, Inches(0.8),
        "GTFS-Realtime is proto2: a field can be **genuinely absent**.\n"
        "`delay = 0` means on time.  Absent means **no information**.\n"
        "Read the natural way, protobuf returns `0` for **both**.",
        size=18, line_spacing=1.4)
    big = panel(s, MARGIN, y + Inches(1.5), W - 2 * MARGIN, Inches(1.85),
                fill=RGBColor(0xFD, 0xF4, 0xE8))
    big.line.color.rgb = RGBColor(0xE8, 0xC4, 0x9A)
    txt(s, MARGIN + Inches(0.4), y + Inches(1.72), Inches(4.2), Inches(1.3),
        "44.2%", size=60, bold=True, color=WARN)
    txt(s, MARGIN + Inches(4.3), y + Inches(1.8), W - 2 * MARGIN - Inches(4.7),
        Inches(1.4),
        "of records carry **no delay value at all** — 54,638 of 123,735.\n"
        "19 of 28 operators never populate it.\n"
        "Written the obvious way, we'd have reported every one as **exactly on time.**",
        size=16, line_spacing=1.4)
    txt(s, MARGIN, y + Inches(3.55), W - 2 * MARGIN, Inches(0.4),
        "Verified against raw protobuf across 408,149 rows — 0 violations.",
        size=14, color=MUTED)
    budget.append(notes(s, 38,
        "BORNA. The single best slide in the deck. A fabricated on-time spike "
        "that looks completely plausible. Then: 'Aatish will show how that data "
        "becomes an answer.' HANDOFF."))

    # ---------------- 6. Architecture ------------------------------------
    s = add_slide(prs)
    y = title(s, "Implemented end-to-end path", speaker="Borna")
    ry = y + Inches(0.15)
    bh = Inches(0.85)
    w1 = Inches(1.85)
    gap = Inches(0.42)
    x = MARGIN
    flow_box(s, x, ry, w1, bh, "511 feed", "protobuf, 120s")
    arrow(s, x + w1 + Inches(0.07), ry + Inches(0.34), gap - Inches(0.14))
    x += w1 + gap
    flow_box(s, x, ry, w1, bh, "poller", "archive first")
    arrow(s, x + w1 + Inches(0.07), ry + Inches(0.34), gap - Inches(0.14))
    x += w1 + gap
    flow_box(s, x, ry, w1, bh, "data/raw/", "system of record",
             fill=RGBColor(0xE8, 0xF1, 0xF7))
    arrow(s, x + w1 + Inches(0.07), ry + Inches(0.34), gap - Inches(0.14))
    x += w1 + gap
    flow_box(s, x, ry, w1, bh, "producer", "validate → publish")
    arrow(s, x + w1 + Inches(0.07), ry + Inches(0.34), gap - Inches(0.14))
    x += w1 + gap
    flow_box(s, x, ry, w1, bh, "Kafka", "4 topics",
             fill=RGBColor(0xE8, 0xF1, 0xF7))

    ry2 = ry + Inches(1.35)
    x = MARGIN + Inches(1.2)
    flow_box(s, x, ry2, Inches(2.3), bh, "arrival resolver", "STOPPED_AT → arrival")
    arrow(s, x + Inches(2.4), ry2 + Inches(0.34), Inches(0.5))
    flow_box(s, x + Inches(3.0), ry2, Inches(2.3), bh, "aggregator", "join predictions")
    arrow(s, x + Inches(5.4), ry2 + Inches(0.34), Inches(0.5))
    flow_box(s, x + Inches(6.0), ry2, Inches(2.0), bh, "outputs/", "metrics",
             fill=RGBColor(0xE8, 0xF1, 0xF7))
    arrow(s, x + Inches(8.1), ry2 + Inches(0.34), Inches(0.5))
    flow_box(s, x + Inches(8.7), ry2, Inches(1.9), bh, "ML model", "corrected ETA")

    txt(s, MARGIN, ry2 + Inches(1.35), W - 2 * MARGIN, Inches(1.2),
        "**Why this shape:** the poller only writes to disk; everything downstream only reads from disk.\n"
        "Live and replay differ **only** in whether new files appear — so the grader runs the real pipeline, not a mock.\n"
        "**Kafka is transport, not storage.** 24h retention on purpose: GTFS-RT has no history endpoint, so the archive is the system of record.",
        size=15, line_spacing=1.4)
    budget.append(notes(s, 62,
        "BORNA. Rubric slide. Trace the path left to right in ONE sentence, "
        "then give the WHY: one code path for live and replay, and Kafka is "
        "transport not storage. Do not read every box. Then hand over: "
        "\'Aatish will show this running on a single vehicle.\'"))

    # ---------------- 7. Concrete example --------------------------------
    s = add_slide(prs)
    y = title(s, "One concrete example: deriving a single arrival",
              speaker="Aatish")
    panel(s, MARGIN, y - Inches(0.1), Inches(7.0), Inches(2.5))
    mono(s, MARGIN + Inches(0.3), y + Inches(0.15), Inches(6.5), Inches(2.2),
         """vehicle 5885, trip SF:12053041_M11

poll 1   IN_TRANSIT_TO  stop_sequence 46
poll 2   **STOPPED_AT     stop_sequence 47**   ← arrival
poll 3   STOPPED_AT     stop_sequence 47    (suppressed)
poll 4   IN_TRANSIT_TO  stop_sequence 48""", size=14)

    xr = MARGIN + Inches(7.4)
    wr = W - MARGIN - xr
    txt(s, xr, y - Inches(0.05), wr, Inches(3.0),
        "The arrival is the **first** sighting,\nnot the last.\n\n"
        "A parked bus reports STOPPED_AT\non every poll — the last one\nwould measure **departure**.\n\n"
        "So the resolver holds per-trip state,\nwhich is **why the Kafka key matters**:\n"
        "all observations of one trip must\narrive in order, in one partition.",
        size=15, line_spacing=1.35)
    txt(s, MARGIN, y + Inches(2.75), W - 2 * MARGIN, Inches(0.9),
        "Every arrival is stamped with **how** it was derived — method, confidence, poll_interval_s.\n"
        "Resolver availability varies by operator (15 of 28), so without provenance "
        "**“agency” silently becomes a proxy for data quality.**",
        size=15, line_spacing=1.4)
    budget.append(notes(s, 45,
        "AATISH. Rubric asks for one concrete example — this is it. The "
        "first-vs-last point is the memorable bit, and it justifies the "
        "partition key decision from Borna's contract slide."))

    # ---------------- 8. Results -----------------------------------------
    s = add_slide(prs)
    y = title(s, "Result: how wrong are Muni's predictions?",
              sub="SF Muni · 25,501 prediction–arrival pairs", speaker="Aatish")
    rows = [
        ("lead time", "n", "bias", "median |err|", "within 60s"),
        ("0–2 min", "509", "−11 s", "13 s", "92.7%"),
        ("2–5 min", "4,603", "−46 s", "50 s", "59.5%"),
        ("5–10 min", "4,897", "−66 s", "70 s", "43.3%"),
        ("10–20 min", "8,919", "−98 s", "102 s", "31.1%"),
        ("20+ min", "6,573", "−142 s", "146 s", "23.0%"),
    ]
    tw = Inches(8.4)
    colw = [Inches(1.9), Inches(1.4), Inches(1.6), Inches(2.0), Inches(1.5)]
    ry = y - Inches(0.1)
    for i, row in enumerate(rows):
        cx = MARGIN
        if i == 0:
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN,
                                    ry + Inches(0.34), tw, Pt(1.5))
            ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
            ln.line.fill.background(); ln.shadow.inherit = False
        for j, cell in enumerate(row):
            txt(s, cx, ry, colw[j], Inches(0.35), cell,
                size=15, bold=(i == 0),
                color=ACCENT if i == 0 else (WARN if (j == 2 and i == 5) else INK))
            cx += colw[j]
        ry += Inches(0.42) if i == 0 else Inches(0.4)

    txt(s, MARGIN, ry + Inches(0.2), W - 2 * MARGIN, Inches(1.3),
        "**Negative bias = the agency predicts EARLIER than reality.** Buses arrive later than promised,\n"
        "and it worsens with horizon. Short-horizon ETAs are excellent; 20-minute ETAs are not.",
        size=16, line_spacing=1.4)
    txt(s, MARGIN + Inches(8.9), y + Inches(0.1), Inches(3.2), Inches(2.6),
        "Credible because the two\nhalves are **independent**:\n\n"
        "prediction ← TripUpdates\nactual ← VehiclePositions\n\n"
        "Deriving arrivals from\npredictions would measure\nthe agency **against itself.**",
        size=14, color=MUTED, line_spacing=1.35)
    budget.append(notes(s, 48,
        "AATISH. Read one row, not five: 'at twenty minutes out, only 23% land "
        "within a minute.' Then the independence point — that is what makes the "
        "number mean anything."))

    # ---------------- 10. Lesson learned ---------------------------------
    s = add_slide(prs)
    y = title(s, "Lesson learned: the wrong answer looked better",
              speaker="Aatish")
    panel(s, MARGIN, y - Inches(0.1), W - 2 * MARGIN, Inches(1.65),
          fill=RGBColor(0xFD, 0xF4, 0xE8))
    txt(s, MARGIN + Inches(0.35), y + Inches(0.12), W - 2 * MARGIN - Inches(0.7),
        Inches(1.3),
        "Day-of-week **improved** test MAE by 10 seconds. We removed it.\n"
        "Monday appeared **only** in the test window — 31% of test rows had a value the model "
        "had never seen, silently inheriting another day's correction through bin placement.",
        size=16, line_spacing=1.4)
    txt(s, MARGIN, y + Inches(1.85), W - 2 * MARGIN, Inches(0.5),
        "**A number you can't account for is not a result.**", size=21, color=WARN)

    txt(s, MARGIN, y + Inches(2.55), W - 2 * MARGIN, Inches(1.3),
        "Same pattern across every real bug we hit: a mislabelled partition column, a duplicate replay that\n"
        "doubled n while leaving every average identical, an API key leaked through an exception message.\n"
        "**None crashed. Several produced better-looking numbers.** Reading the code would have caught none of them —\n"
        "measuring against real data did.",
        size=15, line_spacing=1.4)
    budget.append(notes(s, 42,
        "AATISH. Rubric: one evidence-based lesson. This is the strongest story "
        "in the project — we deleted a feature that improved our score. Land the "
        "bold line and pause."))

    # ---------------- 11. Limitation + next step -------------------------
    s = add_slide(prs)
    y = title(s, "Limitation, and the next step", speaker="Aatish")
    cw = (W - 2 * MARGIN - Inches(0.4)) / 2
    panel(s, MARGIN, y - Inches(0.1), cw, Inches(3.1))
    txt(s, MARGIN + Inches(0.3), y + Inches(0.15), cw - Inches(0.6), Inches(0.4),
        "Limitation — measured", size=16, bold=True, color=ACCENT)
    txt(s, MARGIN + Inches(0.3), y + Inches(0.65), cw - Inches(0.6), Inches(2.3),
        "We capture only **~21% of stop events.**\n"
        "Buses dwell for seconds; we sample every 120s\n(a 60 req/hour quota).\n\n"
        "Worse, it isn't random: capture probability scales\nwith dwell time, so **terminals and layovers are\n"
        "over-represented**, and derived arrivals are\nbiased late.\n\n"
        "**So our headline bias is an upper bound,\nnot a point estimate.**",
        size=14, line_spacing=1.35)

    x2 = MARGIN + cw + Inches(0.4)
    panel(s, x2, y - Inches(0.1), cw, Inches(3.1))
    txt(s, x2 + Inches(0.3), y + Inches(0.15), cw - Inches(0.6), Inches(0.4),
        "Next step — realistic", size=16, bold=True, color=ACCENT)
    txt(s, x2 + Inches(0.3), y + Inches(0.65), cw - Inches(0.6), Inches(2.3),
        "**Request a rate-limit increase from 511.**\n\n"
        "One email. It attacks all three problems at once:\nmore stop events captured, tighter timestamps,\n"
        "and less selection bias toward long dwells.\n\n"
        "Then: GTFS-Static for true on-time performance\nagainst the published schedule, and a geofence\n"
        "resolver for the 13 operators where\nSTOPPED_AT never fires.",
        size=14, line_spacing=1.35)
    budget.append(notes(s, 42,
        "AATISH. Rubric: one limitation + one next step. Be honest about the "
        "upper bound — examiners reward that. The next step is deliberately "
        "small and credible: one email."))

    # ---------------- 12. Close ------------------------------------------
    s = add_slide(prs)
    txt(s, MARGIN, Inches(2.1), W - 2 * MARGIN, Inches(0.8),
        "One command, no API key", size=34, bold=True)
    panel(s, MARGIN, Inches(3.1), Inches(6.2), Inches(1.0))
    mono(s, MARGIN + Inches(0.4), Inches(3.38), Inches(5.4), Inches(0.5),
         "make demo", size=22)
    txt(s, MARGIN + Inches(6.8), Inches(3.15), Inches(5.4), Inches(1.2),
        "Kafka → replay → arrivals → metrics\n→ 6 acceptance checks, in **1 min 59 s**",
        size=16, color=MUTED, line_spacing=1.4)
    txt(s, MARGIN, Inches(4.65), W - 2 * MARGIN, Inches(1.0),
        "83 tests  ·  6/6 acceptance checks  ·  1.67 M records replayed, 0 dead-lettered\n"
        "Verified from a clean extract of the submitted ZIP.",
        size=17, line_spacing=1.45)
    txt(s, MARGIN, Inches(6.15), W - 2 * MARGIN, Inches(0.5),
        "Questions?", size=26, bold=True, color=ACCENT)
    budget.append(notes(s, 12,
        "EITHER. Close on reproducibility, then stop talking. Leave the full "
        "minute for Q&A."))

    # ---------------- A1. APPENDIX: AI element (Q&A backup) --------------------------------------
    s = add_slide(prs)
    y = title(s, "Appendix: bounded AI element", speaker="appendix — Q&A")
    txt(s, MARGIN, y - Inches(0.1), W - 2 * MARGIN, Inches(0.5),
        "Predict the residual  (actual − predicted),  apply it as a correction. "
        "Predicting 0 = trusting the agency.", size=16)
    rows = [
        ("", "MAE", "beats agency?"),
        ("trust the agency", "195.2 s", "—"),
        ("add global mean", "215.1 s", "worse"),
        ("add mean per horizon", "206.7 s", "worse"),
        ("gradient boosting model", "166.6 s", "−14.7%"),
    ]
    ry = y + Inches(0.6)
    colw = [Inches(3.6), Inches(1.7), Inches(1.9)]
    for i, row in enumerate(rows):
        cx = MARGIN
        if i == 0:
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN,
                                    ry + Inches(0.32), Inches(7.2), Pt(1.5))
            ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
            ln.line.fill.background(); ln.shadow.inherit = False
        for j, cell in enumerate(row):
            txt(s, cx, ry, colw[j], Inches(0.35), cell, size=15,
                bold=(i == 0 or i == 4),
                color=ACCENT if i == 0 else (WARN if i == 4 and j == 2 else INK))
            cx += colw[j]
        ry += Inches(0.42) if i == 0 else Inches(0.38)

    xr = MARGIN + Inches(7.8)
    wr = W - MARGIN - xr
    txt(s, xr, y + Inches(0.5), wr, Inches(2.6),
        "**Both naive corrections make it worse.**\n"
        "Residuals are right-skewed, so adding\nthe mean overcorrects the typical case.\n\n"
        "**Leakage guard:** labels come from GPS,\nfeatures from predictions — different\nfeeds, so the label can't contain\nthe feature.",
        size=14, line_spacing=1.35)
    txt(s, MARGIN, y + Inches(3.05), W - 2 * MARGIN, Inches(0.5),
        "Fallback enforced in code: if the model doesn't beat its baseline, serve the agency's ETA unchanged.",
        size=14, color=MUTED)
    notes(s, 0,
          "APPENDIX -- not presented. Use if asked about the AI element.")

    out = Path("Final_Presentation_Karimi_Lobo.pptx")
    prs.save(out)

    total = sum(budget)
    print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
    print(f"speaker-note budget: {total}s of {TOTAL_BUDGET_S}s "
          f"({TOTAL_BUDGET_S - total:+d}s margin)")
    if total > TOTAL_BUDGET_S:
        print("WARNING: over the 7-minute limit -- cut content, the rubric "
              "awards points for staying in time")
    return out


if __name__ == "__main__":
    print(f"wrote {build()}")
