"""Pipeline-organised deck: walk the data from source to answer.

Structure follows the pipeline rather than the rubric categories -- but the
rubric still has to be satisfied, so each stage is placed to cover one:

  slide 2      -> rubric 1  problem, target user, why it matters
  slide 3      -> rubric 3  the architecture chart (whole path, one picture)
  slides 4-5   -> rubric 2  data source, event contract, representative record
  slide 6      -> rubric 3  design choice, explained (Kafka: key, partitions)
  slide 7      -> rubric 3  ONE concrete example (a real vehicle)
  slides 9-10  -> rubric 4  evidence-based lesson + realistic next step

Slot: Thursday 13 August, 6:16-6:24pm. 7 minutes presenting, 1 minute Q&A.
Both speakers must talk. Timing is scored, so budgets are asserted below.

Regenerate:  python scripts/build_deck_pipeline.py
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_deck import (ACCENT, BG, INK, LINE, MARGIN, MUTED, PANEL, W, WARN,
                        add_slide, arrow, flow_box, mono, panel, title, txt)

TOTAL_BUDGET_S = 7 * 60

STAGE = RGBColor(0xE8, 0xF1, 0xF7)   # highlight for the stage under discussion


# ---------------------------------------------------------------------------
# Delivery scripts -- near-verbatim, budgeted at ~150 words per minute.
# ---------------------------------------------------------------------------
SCRIPTS = {
 1: (10, "BORNA", """We measure how late Bay Area transit actually is. The catch that shaped this
whole project is that no agency publishes that number, so we derive it.""",
     "Don't linger — the hook is the next slide."),

 2: (48, "BORNA", """Transit agencies publish two live feeds. TripUpdates gives you predictions: a bus will
reach stop 22 at 15:07. VehiclePositions gives you GPS: bus 4821 is at this
coordinate right now. But notice what neither of them ever says, which is that
the bus actually arrived. That fact is not published by anyone, so it has to be
derived from indirect evidence. And that is what makes this hard rather than
tedious, because delay is actual minus scheduled, the scheduled time is published
and exact, so every bit of error in our derived actual goes straight into the
delay figure with nothing downstream able to catch it. Our user is a rider
deciding whether to trust the estimate in their app.""",
     "Slow down from 'And that is what makes this hard' — it justifies everything after."),

 3: (36, "BORNA", """This is the whole talk in one picture, and we are going to walk it left to
right. 511 gives us protobuf over HTTP. A poller writes those bytes straight to
disk. A producer then decodes them, validates every record, and publishes to
Kafka. A resolver turns GPS observations into arrivals, an aggregator joins those
arrivals back against the original predictions, and that produces our answer.
Finally a model learns the correction. Nine stages in total, and I will take the
first four while Aatish takes the rest.""",
     "Point along the diagram as you go. Don't read every box aloud."),

 4: (40, "BORNA", """Stages one and two. Our source is a single consolidated 511 feed covering all
twenty-eight Bay Area operators in one request, and that matters because our
quota is sixty requests an hour, which means one poll every two minutes. The
poller then has one rule everything depends on: write the raw bytes to disk
before attempting to decode them. A decoding bug is recoverable, because the file is still there
and we can reprocess it. A poll we never took is gone forever, because there is
no history endpoint. So we make the irreversible step the cheap one.""",
     "Pause after the last sentence — it is the line people remember."),

 5: (48, "BORNA", """Stages three and four. The archive on disk is our system of record, and
everything downstream reads from it. On the left is one real record as it enters
Kafka, after validation. The key is service date plus trip ID, and the date is in
there because trip IDs repeat every single day. Now look at arrival uncertainty,
because it is null rather than zero. Protobuf lets a field be genuinely absent, and
absent means something quite different from zero. Zero means the vehicle is
exactly on time. Absent means the agency told us nothing at all. Forty-four
percent of records carry no delay value, and read the obvious way we would have
called every single one of them perfectly on time.""",
     "Land 'perfectly on time', pause, then hand over to Aatish."),

 6: (44, "AATISH", """Stage five is Kafka, and we run four topics. Two of them carry the raw feeds,
predictions and vehicle positions. A third is a dead letter topic for records
that fail validation, so one malformed row degrades the run rather than stopping
it. The fourth carries the arrivals we derive, so anything downstream can consume
them without touching our code.

The three data topics get three partitions each, while the dead letter gets one,
because ordering between unrelated broken records is meaningless. We chose three
deliberately rather than accepting a default, because partition count is
effectively permanent: Kafka routes by hashing the key against it, so adding a
partition later re-routes keys and breaks per-trip ordering retroactively, for
data already written.""",
     "Name what each topic is FOR. The partition point is the design choice being scored."),

 7: (50, "AATISH", """Stage six, and this is the heart of the project. Here is one real vehicle, a
Powell-Mason cable car, across four consecutive polls. The status field comes
straight from the feed, so when Muni's own equipment reports STOPPED_AT at stop
four, we are reading that rather than inferring it. What we derive is the
arrival: the moment that status first appears for that stop. We take the first
sighting rather than the last, because a parked vehicle keeps reporting stopped
on every poll, so the last one would measure departure. That derived sentence,
this trip arrived at stop four at 14:54:31, exists nowhere in the source data,
and producing it requires the resolver to hold per-trip state, which is exactly
why the ordering guarantee matters.""",
     "The read-vs-derived distinction is the point. First-vs-last is now one sentence."),

 8: (43, "AATISH", """Stages seven and eight. Now we join our derived arrivals against what the
agency actually predicted. For that same cable car, two minutes ahead of time
Muni predicted 14:53:54, and it arrived at 14:54:31, so thirty-seven seconds
later than promised. That is one data point, and the pipeline produces
twenty-five thousand of them. Aggregated, the close-in predictions are genuinely
excellent, around thirteen seconds off. But at twenty minutes out only
twenty-three percent land within a minute, and the bias is negative at every
horizon, meaning vehicles consistently arrive later than promised. This is only
credible because the two halves come from completely different feeds.""",
     "Read one row of the table aloud, not five. The independence point is the content."),

 9: (36, "AATISH", """Our best lesson came from a feature we deleted. Day of week improved our test
error by ten seconds, and we removed it anyway. The reason is that Monday
appeared only in our test window, so thirty-one percent of test rows carried a
value the model had never seen in training, so they quietly inherited an adjacent
day's correction through wherever the bin boundaries fell. The gain was an
artifact, not learned structure. Every real bug we hit had that shape: none
crashed, and several produced better-looking numbers than the correct version.""",
     "Pause after 'we removed it anyway' and let the oddity register."),

10: (38, "AATISH", """Our main limitation is coverage, because we catch only about twenty-one
percent of stop events, because vehicles dwell for seconds while we sample every
two minutes. And it is not random, since the chance of catching a stop scales
with dwell time, so terminals are over-represented and our arrivals skew late.
That means the optimism we measured is an upper bound, not a point estimate. So
our next step is deliberately small: one email to 511 requesting a higher rate
limit, which attacks all three problems at once.""",
     "Being honest about the upper bound is worth marks. Don't soften it."),

11: (12, "EITHER", """All of that runs from a single command, with no API key, in about two minutes,
covering eighty-three tests and six acceptance checks. We are happy to take
questions.""",
     "Then stop talking. Leave the full minute for Q&A."),

12: (0, "APPENDIX", """Not presented — this is backup in case we are asked about the AI element.

We predict the residual, which is the actual arrival minus the predicted arrival,
and we apply it as a correction on top of the agency's own estimate. Predicting
zero is identical to trusting the agency, so the baseline comparison is exact
rather than approximate. The interesting result is that both naive bias
corrections perform worse than doing nothing at all, because residuals are
heavily right-skewed, with a mean of plus 213 seconds against a median of plus
111, so adding the mean overcorrects the typical case. The model earns its
fourteen point seven percent improvement by learning a correction that depends on
horizon, hour and route rather than applying a constant.

On leakage: our labels come from GPS positions and our features come from
predictions, so they are different feeds and the label cannot contain the
feature. We also excluded lead time as a feature because it contains the target,
and we split the data strictly forward in time rather than randomly.""",
     "Only if asked."),
}

def pipeline_strip(s, y, highlight=None, h=Inches(0.62), small=True):
    """The nine-stage strip, with the current stage highlighted.

    Repeating it on every slide costs vertical space but keeps the audience
    oriented -- in a seven-minute talk nobody can hold a nine-box diagram in
    their head from slide three onward.
    """
    stages = ["511 feed", "poller", "archive", "producer", "Kafka",
              "resolver", "aggregator", "outputs", "model"]
    n = len(stages)
    gap = Inches(0.12)
    bw = (W - 2 * MARGIN - gap * (n - 1)) / n
    for i, name in enumerate(stages):
        x = MARGIN + i * (bw + gap)
        on = highlight is not None and i in highlight
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, h)
        b.fill.solid()
        b.fill.fore_color.rgb = STAGE if on else RGBColor(0xF7, 0xF9, 0xFB)
        b.line.color.rgb = ACCENT if on else LINE
        b.line.width = Pt(1.5 if on else 0.75)
        b.shadow.inherit = False
        b.adjustments[0] = 0.18
        p = b.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i+1}. {name}"
        r.font.size = Pt(10.5 if small else 12)
        r.font.bold = on
        r.font.color.rgb = INK if on else MUTED
        r.font.name = "Calibri"
        if i < n - 1:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                   x + bw + Inches(0.02), y + h / 2 - Inches(0.045),
                                   Inches(0.08), Inches(0.09))
            a.rotation = 0
            a.fill.solid(); a.fill.fore_color.rgb = LINE
            a.line.fill.background(); a.shadow.inherit = False
    return y + h


def apply_scripts(prs):
    total = 0
    for idx, slide in enumerate(prs.slides, 1):
        if idx not in SCRIPTS:
            continue
        secs, speaker, script, delivery = SCRIPTS[idx]
        total += secs
        words = len(script.split())
        wpm = round(words / (secs / 60)) if secs else 0
        slide.notes_slide.notes_text_frame.text = (
            f"[{secs}s  |  {speaker}  |  ~{words} words @ {wpm} wpm]\n"
            f"{'=' * 62}\n\n{script.strip()}\n\n"
            f"{'-' * 62}\nDELIVERY: {delivery}"
        )
    return total


def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, Inches(7.5)

    # ---- 1. Title -------------------------------------------------------
    s = add_slide(prs)
    txt(s, MARGIN, Inches(2.2), W - 2 * MARGIN, Inches(1.0),
        "Real-Time Transit Reliability Lakehouse", size=42, bold=True)
    txt(s, MARGIN, Inches(3.2), W - 2 * MARGIN, Inches(0.8),
        "Deriving arrival times that no agency publishes —\n"
        "and grading the predictions riders actually see",
        size=20, color=MUTED, line_spacing=1.3)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.3), Inches(2.0), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT
    r.line.fill.background(); r.shadow.inherit = False
    txt(s, MARGIN, Inches(4.7), W - 2 * MARGIN, Inches(0.9),
        "Borna Karimi  ·  Aatish Lobo\nMSDS 682 — Data Stream Processing",
        size=17, line_spacing=1.35)
    pipeline_strip(s, Inches(6.15))
    txt(s, MARGIN, Inches(6.95), W - 2 * MARGIN, Inches(0.3),
        "Data provided by 511.org (Metropolitan Transportation Commission)",
        size=10, color=MUTED)

    # ---- 2. The problem -------------------------------------------------
    s = add_slide(prs)
    y = title(s, "Nobody publishes when a bus arrives", speaker="Borna")
    cw = (W - 2 * MARGIN - Inches(0.35)) / 2
    panel(s, MARGIN, y - Inches(0.15), cw, Inches(1.35))
    txt(s, MARGIN + Inches(0.25), y + Inches(0.05), cw - Inches(0.5), Inches(1.0),
        "TripUpdates — **predictions**\n“will reach stop 22 at 15:07”",
        size=16, line_spacing=1.35)
    panel(s, MARGIN + cw + Inches(0.35), y - Inches(0.15), cw, Inches(1.35))
    txt(s, MARGIN + cw + Inches(0.6), y + Inches(0.05), cw - Inches(0.5), Inches(1.0),
        "VehiclePositions — **GPS pings**\n“bus 4821 is here, now”",
        size=16, line_spacing=1.35)
    txt(s, MARGIN, y + Inches(1.45), W - 2 * MARGIN, Inches(0.5),
        "Neither ever says: “bus 4821 **arrived** at stop 22 at 15:09.”", size=21)
    txt(s, MARGIN, y + Inches(2.1), W - 2 * MARGIN, Inches(1.1),
        "delay = actual − scheduled,  and **scheduled is exact** —\n"
        "so every error in our derivation lands **straight in the answer**.",
        size=18, line_spacing=1.45)
    txt(s, MARGIN, y + Inches(3.15), W - 2 * MARGIN, Inches(0.4),
        "**User:** a rider deciding whether to trust the ETA in their app.",
        size=15, color=MUTED)
    pipeline_strip(s, Inches(6.5))

    # ---- 3. The pipeline ------------------------------------------------
    s = add_slide(prs)
    y = title(s, "The whole pipeline", sub="We'll walk it left to right",
              speaker="Borna")
    bh = Inches(0.8); w1 = Inches(1.72); gap = Inches(0.30)
    row = [("1. 511 feed","protobuf"),("2. poller","archive first"),
           ("3. data/raw/","system of record"),("4. producer","validate"),
           ("5. Kafka","4 topics")]
    x = MARGIN
    for lbl, sub in row:
        flow_box(s, x, y - Inches(0.1), w1, bh, lbl, sub,
                 fill=STAGE if "raw" in lbl or "Kafka" in lbl else PANEL)
        if x + w1 < W - MARGIN - w1:
            arrow(s, x + w1 + Inches(0.04), y + Inches(0.22), gap - Inches(0.08))
        x += w1 + gap
    row2 = [("6. resolver","GPS → arrivals"),("7. aggregator","join predictions"),
            ("8. outputs/","the metrics"),("9. model","learn correction")]
    x = MARGIN + Inches(1.0)
    y2 = y + Inches(1.25)
    for lbl, sub in row2:
        flow_box(s, x, y2, Inches(2.2), bh, lbl, sub,
                 fill=STAGE if "outputs" in lbl else PANEL)
        if lbl != row2[-1][0]:
            arrow(s, x + Inches(2.24), y2 + Inches(0.32), Inches(0.3))
        x += Inches(2.5)
    txt(s, MARGIN, y2 + Inches(1.25), W - 2 * MARGIN, Inches(0.9),
        "**Borna:** stages 1–4, getting the data in and validated.\n"
        "**Aatish:** stages 5–9, turning it into an answer.",
        size=17, line_spacing=1.4)

    # ---- 4. Stages 1-2 --------------------------------------------------
    s = add_slide(prs)
    y = title(s, "1–2 · The source, and one rule", speaker="Borna")
    panel(s, MARGIN, y - Inches(0.15), Inches(5.6), Inches(1.9))
    txt(s, MARGIN + Inches(0.3), y + Inches(0.05), Inches(5.0), Inches(1.6),
        "**511.org** — one regional feed\n**28 operators**, GTFS-Realtime protobuf\n"
        "**60 requests/hour** → one poll per 120 s",
        size=16, line_spacing=1.5)
    xr = MARGIN + Inches(6.0)
    txt(s, xr, y - Inches(0.1), W - MARGIN - xr, Inches(2.0),
        "The quota is not just a limit —\nit is the **precision limit** on every\n"
        "arrival we derive. So it travels\nwith every row as `poll_interval_s`.",
        size=15, line_spacing=1.4)
    big = panel(s, MARGIN, y + Inches(2.1), W - 2 * MARGIN, Inches(1.3),
                fill=RGBColor(0xFD, 0xF4, 0xE8))
    big.line.color.rgb = RGBColor(0xE8, 0xC4, 0x9A)
    txt(s, MARGIN + Inches(0.35), y + Inches(2.3), W - 2 * MARGIN - Inches(0.7),
        Inches(1.0),
        "**Write raw bytes to disk BEFORE decoding.**  A decode bug is recoverable — "
        "reprocess the file.\nA poll never taken is gone forever. "
        "**Make the irreversible step the cheap one.**",
        size=16, line_spacing=1.4)
    pipeline_strip(s, Inches(6.5), highlight={0, 1})

    # ---- 5. Stages 3-4 --------------------------------------------------
    s = add_slide(prs)
    y = title(s, "3–4 · Archive, then a validated event", speaker="Borna")
    panel(s, MARGIN, y - Inches(0.2), Inches(6.5), Inches(3.4))
    mono(s, MARGIN + Inches(0.28), y - Inches(0.02), Inches(6.0), Inches(3.1),
         """key = "20260806:SF:12086593_M11"

{
  "envelope": {
    "poll_interval_s": 120,
    "ingest_ts": "...",       // when WE got it
  },
  "trip": {
    "trip_id":      "SF:12086593_M11",
    "route_id":     "SF:PM",
    "service_date": "20260806"
  },
  "stop_sequence":       4,
  "arrival_time_epoch":  1786053234,
  "arrival_uncertainty": **null**
}""", size=12.5)
    xr = MARGIN + Inches(6.9)
    wr = W - MARGIN - xr
    txt(s, xr, y - Inches(0.15), wr, Inches(1.4),
        "**Key = service_date : trip_id**\ntrip IDs repeat every day, and this\n"
        "keeps one trip in one partition,\nin order",
        size=14, line_spacing=1.35)
    b = panel(s, xr, y + Inches(1.3), wr, Inches(1.95),
              fill=RGBColor(0xFD, 0xF4, 0xE8))
    b.line.color.rgb = RGBColor(0xE8, 0xC4, 0x9A)
    txt(s, xr + Inches(0.25), y + Inches(1.5), wr - Inches(0.5), Inches(1.6),
        "**null ≠ 0**\n\n`0` = exactly on time\nabsent = no information\n\n"
        "**44.2%** of records have\nno delay value at all",
        size=14, color=WARN, line_spacing=1.35)
    pipeline_strip(s, Inches(6.5), highlight={2, 3})

    # ---- 6. Stage 5: Kafka ----------------------------------------------
    s = add_slide(prs)
    y = title(s, "5 · Kafka — and why these settings", speaker="Aatish")
    rows = [("gtfsrt.trip_updates.v1", "3", "24 h"),
            ("gtfsrt.vehicle_positions.v1", "3", "24 h"),
            ("gtfsrt.dead_letter.v1", "1", "7 d"),
            ("transit.arrival_events.v1", "3", "7 d")]
    txt(s, MARGIN, y - Inches(0.2), Inches(4.0), Inches(0.3), "topic",
        size=13, bold=True, color=ACCENT)
    txt(s, MARGIN + Inches(4.0), y - Inches(0.2), Inches(1.2), Inches(0.3),
        "parts", size=13, bold=True, color=ACCENT)
    txt(s, MARGIN + Inches(5.1), y - Inches(0.2), Inches(1.2), Inches(0.3),
        "retention", size=13, bold=True, color=ACCENT)
    ry = y + Inches(0.18)
    for t, p_, r_ in rows:
        txt(s, MARGIN, ry, Inches(4.0), Inches(0.3), t, size=13)
        txt(s, MARGIN + Inches(4.0), ry, Inches(1.2), Inches(0.3), p_, size=13)
        txt(s, MARGIN + Inches(5.1), ry, Inches(1.2), Inches(0.3), r_, size=13)
        ry += Inches(0.34)

    xr = MARGIN + Inches(6.6)
    wr = W - MARGIN - xr
    txt(s, xr, y - Inches(0.2), wr, Inches(3.2),
        "**retention = 24 h, on purpose**\nKafka is transport. The archive is\nstorage.\n\n"
        "**delete, never compact**\nCompaction keeps only the newest\nmessage per key — it would discard\n"
        "the state transitions we read.\n\n"
        "**3 partitions, chosen not defaulted**\nCount is effectively permanent:\n"
        "changing it re-routes keys and\nbreaks ordering retroactively.",
        size=13.5, line_spacing=1.32)
    pipeline_strip(s, Inches(6.5), highlight={4})

    # ---- 7. Stage 6: the resolver ---------------------------------------
    s = add_slide(prs)
    y = title(s, "6 · The resolver — one real vehicle", speaker="Aatish")
    panel(s, MARGIN, y - Inches(0.2), Inches(6.9), Inches(2.35))
    mono(s, MARGIN + Inches(0.3), y - Inches(0.02), Inches(6.4), Inches(2.0),
         """Powell–Mason cable car, trip SF:12086593_M11

14:52:43   IN_TRANSIT_TO   stop_sequence 4
14:54:31   **STOPPED_AT      stop_sequence 4     <- ARRIVED**
14:56:50   STOPPED_AT      stop_sequence 6
14:59:55   STOPPED_AT      stop_sequence 10""", size=13.5)
    xr = MARGIN + Inches(7.3)
    wr = W - MARGIN - xr
    txt(s, xr, y - Inches(0.2), wr, Inches(2.6),
        "**READ from the feed:**\n`current_status` is a published\nGTFS-RT field — we don't infer it.\n\n"
        "**DERIVED by us:**\nthe arrival — the moment that\nstatus *first* appears for a stop.\n\n"
        "First, not last: a parked vehicle\nreports it every poll, so the last\nwould be **departure**.",
        size=14, line_spacing=1.32)
    txt(s, MARGIN, y + Inches(2.4), W - 2 * MARGIN, Inches(0.55),
        "Every arrival is stamped with **how** it was derived: method, confidence, "
        "poll_interval_s.  Only 15 of 28 operators support this.",
        size=14, color=MUTED, line_spacing=1.35)
    pipeline_strip(s, Inches(6.5), highlight={5})

    # ---- 8. Stages 7-8: join + results ----------------------------------
    s = add_slide(prs)
    y = title(s, "7–8 · Join, and the answer", speaker="Aatish")
    panel(s, MARGIN, y - Inches(0.2), Inches(5.0), Inches(1.25))
    txt(s, MARGIN + Inches(0.25), y - Inches(0.02), Inches(4.5), Inches(1.0),
        "Our cable car:\n**predicted 14:53:54 · arrived 14:54:31 → 37 s late**",
        size=14.5, line_spacing=1.4)
    txt(s, MARGIN + Inches(5.3), y + Inches(0.05), Inches(3.0), Inches(0.9),
        "one of **25,501** pairs", size=15, color=MUTED)

    rows = [("lead time","n","bias","med |err|","<60s"),
            ("0–2 min","509","−11 s","13 s","92.7%"),
            ("2–5 min","4,603","−46 s","50 s","59.5%"),
            ("5–10 min","4,897","−66 s","70 s","43.3%"),
            ("10–20 min","8,919","−98 s","102 s","31.1%"),
            ("20+ min","6,573","−142 s","146 s","23.0%")]
    colw=[Inches(1.7),Inches(1.3),Inches(1.4),Inches(1.6),Inches(1.3)]
    ry = y + Inches(1.35)
    for i,row in enumerate(rows):
        cx = MARGIN
        if i==0:
            ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,MARGIN,ry+Inches(0.3),Inches(7.3),Pt(1.5))
            ln.fill.solid(); ln.fill.fore_color.rgb=ACCENT
            ln.line.fill.background(); ln.shadow.inherit=False
        for j,cell in enumerate(row):
            txt(s,cx,ry,colw[j],Inches(0.32),cell,size=14,bold=(i==0),
                color=ACCENT if i==0 else (WARN if j==2 and i==5 else INK))
            cx+=colw[j]
        ry += Inches(0.38) if i==0 else Inches(0.33)
    txt(s, MARGIN + Inches(7.8), y + Inches(1.4), W - MARGIN - (MARGIN+Inches(7.8)),
        Inches(2.2),
        "Credible only because the\ntwo halves are **independent**:\n\n"
        "predicted ← TripUpdates\nactual ← VehiclePositions\n\n"
        "Deriving one from the other\nmeasures the agency\n**against itself.**",
        size=13.5, color=MUTED, line_spacing=1.32)
    pipeline_strip(s, Inches(6.5), highlight={6, 7})

    # ---- 9. Lesson ------------------------------------------------------
    s = add_slide(prs)
    y = title(s, "Lesson: the wrong answer looked better", speaker="Aatish")
    b = panel(s, MARGIN, y - Inches(0.2), W - 2*MARGIN, Inches(1.5),
              fill=RGBColor(0xFD,0xF4,0xE8))
    b.line.color.rgb = RGBColor(0xE8,0xC4,0x9A)
    txt(s, MARGIN+Inches(0.35), y, W-2*MARGIN-Inches(0.7), Inches(1.2),
        "Day-of-week **improved** test error by 10 seconds. We deleted it.\n"
        "Monday appeared **only** in the test window — 31% of test rows carried a value the "
        "model had never seen, inheriting another day's correction through bin placement.",
        size=15.5, line_spacing=1.4)
    txt(s, MARGIN, y + Inches(1.6), W - 2*MARGIN, Inches(0.5),
        "**A number you can't account for is not a result.**", size=21, color=WARN)
    txt(s, MARGIN, y + Inches(2.3), W - 2*MARGIN, Inches(1.1),
        "Same shape as every real bug we hit — a mislabelled partition column, a duplicate replay that\n"
        "doubled n while leaving every average identical. **None crashed. Several looked better.**",
        size=15, line_spacing=1.4)
    pipeline_strip(s, Inches(6.5))

    # ---- 10. Limitation + next ------------------------------------------
    s = add_slide(prs)
    y = title(s, "Limitation, and the next step", speaker="Aatish")
    cw = (W - 2*MARGIN - Inches(0.4)) / 2
    panel(s, MARGIN, y - Inches(0.2), cw, Inches(2.9))
    txt(s, MARGIN+Inches(0.3), y, cw-Inches(0.6), Inches(0.35),
        "Limitation — measured", size=15, bold=True, color=ACCENT)
    txt(s, MARGIN+Inches(0.3), y+Inches(0.5), cw-Inches(0.6), Inches(2.2),
        "We capture only **~21% of stop events**.\nVehicles dwell for seconds; we sample\nevery 120 s.\n\n"
        "And it isn't random — capture scales with\ndwell, so **terminals are over-represented**\n"
        "and our arrivals skew late.\n\n**So the bias we report is an upper bound.**",
        size=13.5, line_spacing=1.35)
    x2 = MARGIN + cw + Inches(0.4)
    panel(s, x2, y - Inches(0.2), cw, Inches(2.9))
    txt(s, x2+Inches(0.3), y, cw-Inches(0.6), Inches(0.35),
        "Next step — realistic", size=15, bold=True, color=ACCENT)
    txt(s, x2+Inches(0.3), y+Inches(0.5), cw-Inches(0.6), Inches(2.2),
        "**Email 511 for a rate-limit increase.**\n\nOne email. It attacks all three at once:\n"
        "more events captured, tighter timestamps,\nless selection bias.\n\n"
        "Then: GTFS-Static for true on-time\nperformance, and a geofence resolver for\nthe 13 operators where this method\ncan't fire.",
        size=13.5, line_spacing=1.35)
    pipeline_strip(s, Inches(6.5))

    # ---- 11. Close ------------------------------------------------------
    s = add_slide(prs)
    txt(s, MARGIN, Inches(1.9), W - 2*MARGIN, Inches(0.8),
        "One command, no API key", size=34, bold=True)
    panel(s, MARGIN, Inches(2.9), Inches(5.6), Inches(0.95))
    mono(s, MARGIN+Inches(0.4), Inches(3.15), Inches(5.0), Inches(0.5),
         "make demo", size=22)
    txt(s, MARGIN+Inches(6.2), Inches(2.95), Inches(6.0), Inches(1.2),
        "Kafka → replay → arrivals → metrics\n→ 6 acceptance checks, in **~2 minutes**",
        size=16, color=MUTED, line_spacing=1.4)
    txt(s, MARGIN, Inches(4.3), W - 2*MARGIN, Inches(1.0),
        "83 tests  ·  6/6 acceptance checks  ·  1.67 M records replayed, 0 dead-lettered\n"
        "Verified from a clean extract of the submitted ZIP.",
        size=16, line_spacing=1.45)
    txt(s, MARGIN, Inches(5.5), W - 2*MARGIN, Inches(0.5),
        "Questions?", size=26, bold=True, color=ACCENT)
    pipeline_strip(s, Inches(6.5))

    # ---- 12. Appendix ---------------------------------------------------
    s = add_slide(prs)
    y = title(s, "Appendix: the bounded AI element", speaker="appendix — Q&A")
    txt(s, MARGIN, y - Inches(0.15), W - 2*MARGIN, Inches(0.5),
        "Predict the residual (actual − predicted), apply as a correction. "
        "Predicting 0 = trusting the agency.", size=15)
    rows=[("","MAE","vs agency"),("trust the agency","195.2 s","—"),
          ("add global mean","215.1 s","worse"),("add mean per horizon","206.7 s","worse"),
          ("gradient boosting model","166.6 s","−14.7%")]
    ry = y + Inches(0.55); colw=[Inches(3.4),Inches(1.6),Inches(1.7)]
    for i,row in enumerate(rows):
        cx=MARGIN
        if i==0:
            ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,MARGIN,ry+Inches(0.3),Inches(6.7),Pt(1.5))
            ln.fill.solid(); ln.fill.fore_color.rgb=ACCENT
            ln.line.fill.background(); ln.shadow.inherit=False
        for j,cell in enumerate(row):
            txt(s,cx,ry,colw[j],Inches(0.32),cell,size=14,bold=(i in (0,4)),
                color=ACCENT if i==0 else (WARN if i==4 and j==2 else INK))
            cx+=colw[j]
        ry += Inches(0.38) if i==0 else Inches(0.34)
    txt(s, MARGIN+Inches(7.3), y+Inches(0.5), W-MARGIN-(MARGIN+Inches(7.3)), Inches(2.6),
        "**Both naive corrections are worse.**\nResiduals are right-skewed, so adding\nthe mean overcorrects.\n\n"
        "**Leakage guard:** labels from GPS,\nfeatures from predictions — different\nfeeds.\n\n"
        "**Fallback in code:** if the model\ndoesn't beat baseline, serve the\nagency ETA unchanged.",
        size=13.5, line_spacing=1.32)

    spoken = apply_scripts(prs)
    out = Path("Final_Presentation_Pipeline_Karimi_Lobo.pptx")
    prs.save(out)
    print(f"slides: {len(prs.slides._sldIdLst)}")
    print(f"spoken: {spoken}s / {TOTAL_BUDGET_S}s ({TOTAL_BUDGET_S - spoken:+d}s)")
    if spoken > TOTAL_BUDGET_S:
        print("WARNING: over time -- the rubric scores staying within it")
    return out


if __name__ == "__main__":
    print(f"wrote {build()}")
